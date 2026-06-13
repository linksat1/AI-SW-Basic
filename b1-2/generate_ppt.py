#!/usr/bin/env python3
"""B1-2 미션 발표자료(PPTX) 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상/크기 정의 ──────────────────────────────────────────────
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
ACCENT  = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
GRAY_BG = RGBColor(0xF2, 0xF2, 0xF2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x33, 0x33, 0x33)
LIGHTBLUE = RGBColor(0xD9, 0xE2, 0xF3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── 헬퍼 함수 ────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)


def add_title_bar(slide, title, accent=PRIMARY, size=30):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return bar


def add_bullets(slide, items, left, top, width, height, base_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = item[0], item[1]
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else DARK
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(max(base_size - level * 2, 12))
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_code_box(slide, lines, left, top, width, height, font_size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_BG
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = "Menlo"
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK
    return box


def add_table(slide, data, left, top, width, height, col_widths=None, font_size=14):
    rows, cols = len(data), len(data[0])
    gshape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF7, 0xF9, 0xFC)
    return table


def add_page_number(slide, num):
    box = slide.shapes.add_textbox(Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT


def case_evidence_slide(num, title, info_bullets, code_lines, extra_bullets=None):
    s = new_slide()
    add_title_bar(s, title, accent=ACCENT)
    add_bullets(s, info_bullets, Inches(0.5), Inches(1.25), Inches(12.3), Inches(1.3), base_size=18)
    code_top = Inches(2.4)
    if extra_bullets:
        add_code_box(s, code_lines, Inches(0.5), code_top, Inches(12.3), Inches(3.4), font_size=13)
        add_bullets(s, extra_bullets, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2), base_size=16)
    else:
        add_code_box(s, code_lines, Inches(0.5), code_top, Inches(12.3), Inches(4.6), font_size=13)
    add_page_number(s, num)
    return s


def case_cause_slide(num, title, cause_bullets, ba_rows, root_fix):
    s = new_slide()
    add_title_bar(s, title, accent=GREEN)
    add_bullets(s, cause_bullets, Inches(0.5), Inches(1.25), Inches(12.3), Inches(2.0), base_size=18)
    add_table(s, ba_rows, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.6),
              col_widths=[Inches(2.0), Inches(4.0), Inches(6.3)], font_size=15)
    add_bullets(s, [("✅ 근본 해결: " + root_fix, 0, True, GREEN)],
                Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), base_size=18)
    add_page_number(s, num)
    return s


def qa_slide(num, title, qa_pairs, accent=PRIMARY):
    s = new_slide()
    add_title_bar(s, title, accent=accent)
    items = []
    for q, a in qa_pairs:
        items.append(("Q. " + q, 0, True, ACCENT))
        items.append(("A. " + a, 0, False, DARK))
        items.append(("", 0))
    add_bullets(s, items, Inches(0.5), Inches(1.25), Inches(12.3), Inches(6.0), base_size=17)
    add_page_number(s, num)
    return s


# ══════════════════════════════════════════════════════════════
# Slide 1 — 표지
# ══════════════════════════════════════════════════════════════
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY
bg.line.fill.background()

box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6))
p = box.text_frame.paragraphs[0]
p.text = "B1-2 : 리눅스 프로세스 및 시스템 리소스 트러블슈팅"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
box.text_frame.word_wrap = True

box2 = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.9))
p2 = box2.text_frame.paragraphs[0]
p2.text = "Memory Leak · CPU Spike · Deadlock — 재현, 분석, 보고"
p2.font.size = Pt(22)
p2.font.color.rgb = LIGHTBLUE

box3 = s.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9))
tf3 = box3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "환경: macOS + OrbStack + Ubuntu 22.04 (agent-admin 계정)"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(0xBF, 0xCC, 0xE6)
p4 = tf3.add_paragraph()
p4.text = "실습 도구: agent-leak-sim.py (Python 시뮬레이터, sudo 불필요)"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0xBF, 0xCC, 0xE6)


# ══════════════════════════════════════════════════════════════
# Slide 2 — 미션 개요
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "미션 개요 — 목표와 3가지 장애 유형")
add_bullets(s, [
    ("한 줄 요약: 실제 서버 장애 3가지를 직접 재현하고, 로그를 증거로 원인을 분석하여", 0, True),
    ("GitHub Issue 형태의 기술 보고서를 작성한다", 1, True),
    ("핵심 전환: “장애 = 감으로 추측하는 것”  →  “장애 = 데이터로 증명하고 기록하는 것”", 0, True, ACCENT),
], Inches(0.5), Inches(1.25), Inches(12.3), Inches(1.6), base_size=18)

add_table(s, [
    ["장애 유형", "현상", "원인"],
    ["Memory Leak (OOM)", "프로세스가 갑자기 종료됨", "메모리를 계속 할당만 하고 해제하지 않음"],
    ["CPU Spike", "CPU 사용률이 급상승 후 프로세스 종료", "특정 연산이 CPU를 독점"],
    ["Deadlock", "프로세스가 살아있으나 완전히 멈춤", "두 스레드가 서로의 자원을 무한 대기"],
], Inches(0.5), Inches(3.0), Inches(12.3), Inches(2.4),
   col_widths=[Inches(2.6), Inches(4.8), Inches(4.9)], font_size=18)

add_bullets(s, [
    ("미션 완료 후 할 수 있는 것: 메모리/CPU 이상 패턴을 로그로 식별, Deadlock을 시스템 도구로 진단,", 0),
    ("그리고 이를 GitHub Issue 형태의 기술 리포트로 작성하여 팀과 공유", 1),
], Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2), base_size=16)
add_page_number(s, 2)


# ══════════════════════════════════════════════════════════════
# Slide 3 — 환경 설정 ①
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "환경 설정 ① — 기본 환경 & B1-2 추가 환경변수")
add_bullets(s, [
    ("실행 환경: macOS + OrbStack + Ubuntu 22.04 (b1-lab VM, agent-admin 계정)", 0),
    ("B1-1 기반 환경 재사용: AGENT_HOME / AGENT_PORT 환경변수, /home/agent-admin/agent-app/ 디렉토리,", 0),
    ("monitor.sh (매분 자동 실행되는 리소스 모니터)", 1),
    ("환경 복구 절차:  chmod +x 'b1-1 setup.sh'  →  sudo bash 'b1-1 setup.sh'", 0),
], Inches(0.5), Inches(1.25), Inches(12.3), Inches(2.2), base_size=18)

add_table(s, [
    ["환경변수 (.bashrc)", "기본값", "의미 / 허용 범위"],
    ["MEMORY_LIMIT", "256", "메모리 상한선 (MB, 50~512)"],
    ["CPU_MAX_OCCUPY", "50", "CPU 최대 점유율 (%, 10~100)"],
    ["MULTI_THREAD_ENABLE", "true", "멀티스레드(Deadlock 시나리오) 사용 여부 (true/false)"],
], Inches(0.5), Inches(3.5), Inches(12.3), Inches(2.0),
   col_widths=[Inches(3.2), Inches(2.0), Inches(7.1)], font_size=17)

add_bullets(s, [
    ("적용 방법: .bashrc에 export 구문 추가 → source ~/.bashrc → echo $VAR로 값 확인", 0, True, GREEN),
], Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0), base_size=18)
add_page_number(s, 3)


# ══════════════════════════════════════════════════════════════
# Slide 4 — 환경 설정 ②
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "환경 설정 ② — 원본 바이너리 대안: agent-leak-sim.py")
add_bullets(s, [
    ("문제: 원본 agent-leak-app은 Linux 바이너리이며, 실습 PC는 sudo 권한이 제한되어", 0, True, ACCENT),
    ("/var/log/agent-app/ 같은 시스템 경로에 자유롭게 쓸 수 없음", 1, False, ACCENT),
    ("해결: 동일한 로그 패턴(형식·문구)을 그대로 출력하는 Python 시뮬레이터를 직접 제작", 0, True, GREEN),
    ("실제 메모리 할당(bytearray 25MB), 실제 CPU 점유(math.sqrt busy-loop),", 1),
    ("실제 threading.Lock 기반 Deadlock 구현 — 모두 sudo 없이 동작", 1),
    ("로그 경로 자동 결정: /var/log/agent-app/ 쓰기 가능하면 사용, 아니면 ./logs/ 로 대체", 1),
], Inches(0.5), Inches(1.25), Inches(12.3), Inches(2.7), base_size=18)

add_table(s, [
    ["케이스", "MEMORY_LIMIT", "CPU_MAX_OCCUPY", "MULTI_THREAD_ENABLE"],
    ["1 — Memory Leak", "256 MB", "100 %", "false"],
    ["2 — CPU Spike", "9999 MB", "50 %", "false"],
    ["3 — Deadlock", "9999 MB", "100 %", "true"],
], Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.8),
   col_widths=[Inches(3.6), Inches(2.9), Inches(2.9), Inches(2.9)], font_size=17)

add_bullets(s, [
    ("실행: python3 agent-leak-sim.py  →  메뉴에서 1/2/3 선택 시 위 프리셋이 자동 적용됨", 0),
], Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), base_size=16)
add_page_number(s, 4)


# ══════════════════════════════════════════════════════════════
# Slide 5 — Case 1 현상 & 증거
# ══════════════════════════════════════════════════════════════
case_evidence_slide(
    5, "Case 1 — Memory Leak: 현상 & 증거",
    info_bullets=[
        ("현상: 메모리가 25MB씩 지속 증가하여 MEMORY_LIMIT(256MB)을 초과", 0, True),
        ("조건: MEMORY_LIMIT=256 / CPU_MAX_OCCUPY=100 / MULTI_THREAD_ENABLE=false", 0),
    ],
    code_lines=[
        "[agent_app.log 발췌]",
        "09:34:48 [INFO] [Memory] Increasing... (+25 MB) Total: 25 MB",
        "09:34:50 [INFO] [Memory] Increasing... (+25 MB) Total: 50 MB",
        "...",
        "09:35:32 [INFO] [Memory] Increasing... (+25 MB) Total: 250 MB",
        "09:35:38 [INFO] [Memory] Increasing... (+25 MB) Total: 275 MB  <- 한계(256MB) 초과",
        "",
        "[CRITICAL] [MemoryGuard] Memory limit exceeded (275MB >= 256MB)",
        "[CRITICAL] [MemoryGuard] Self-terminating process XXXX to prevent system instability.",
        ">>> [SYSTEM] SELF-TERMINATED (Memory Limit Exceeded) <<<",
    ],
    extra_bullets=[
        ("메모리 상승 속도: 0MB → 275MB 약 50초, +25MB/단계로 일정하게 증가", 0),
    ],
)


# ══════════════════════════════════════════════════════════════
# Slide 6 — Case 1 원인 & Before/After
# ══════════════════════════════════════════════════════════════
case_cause_slide(
    6, "Case 1 — Memory Leak: 원인 분석 & 조치",
    cause_bullets=[
        ("앱 내부에서 메모리를 할당 후 해제하지 않는 메모리 누수(Memory Leak) 발생", 0, True),
        ("25MB 단위로 계속 누적 → MEMORY_LIMIT(256MB) 초과 시 MemoryGuard가 SIGKILL로 자가 종료", 0),
        ("정상 동작: 할당 → 사용 → 해제 (일정 유지)   |   누수 동작: 할당 → 사용 → [해제 안 함] → 계속 증가", 0),
    ],
    ba_rows=[
        ["구분", "설정", "결과"],
        ["Before", "MEMORY_LIMIT = 256 MB", "약 50초 후 한계 도달 → SELF-TERMINATED"],
        ["After", "MEMORY_LIMIT = 512 MB", "더 오랜 시간 생존 확인 (한계치 상향만으로는 결국 재발)"],
    ],
    root_fix="소스코드에서 사용 후 불필요한 데이터를 주기적으로 해제(del / free)하는 리팩토링",
)


# ══════════════════════════════════════════════════════════════
# Slide 7 — Case 2 현상 & 증거
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "Case 2 — CPU Spike: 현상 & 증거", accent=ACCENT)
add_bullets(s, [
    ("현상: CPU 레벨이 1→10까지 램프업, Lv≥8 구간에서 실제 CPU 사용률 99.7% > 50%(임계치)", 0, True),
    ("조건: CPU_MAX_OCCUPY=50 / MEMORY_LIMIT=9999 / MULTI_THREAD_ENABLE=false", 0),
], Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.9), base_size=18)

add_code_box(s, [
    "[agent_app.log 발췌]",
    "17:37:50 [INFO] --- Step Info: Mode=UP, CPU Lv=8, Mem=200MB ---",
    "17:37:50 [INFO] [CPU] Occupy core for 4.0s (Level 8)",
    "17:37:50 [WARNING]  [Watchdog] CPU usage spike detected: 99.7% > 50%",
    "17:37:50 [CRITICAL] [Watchdog] INITIATING EMERGENCY ABORT (SIGTERM)",
    ">>> [SYSTEM] WATCHDOG: PROCESS TERMINATED <<<",
], Inches(0.5), Inches(2.15), Inches(12.3), Inches(1.85), font_size=13)

add_bullets(s, [
    ("CPU Lv=10 도달 패턴이 약 1분 42초 주기로 반복 관측됨 — 아래 4회 관측 시각으로 증명", 0, True, ACCENT),
], Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.45), base_size=17)

add_table(s, [
    ["발생 시각 (agent_app.log)", "CPU Lv", "Mem", "이전 관측 대비 간격"],
    ["15:38:19", "10", "225MB", "—"],
    ["15:40:01", "10", "225MB", "1분 42초 (102초)"],
    ["15:41:43", "10", "225MB", "1분 42초 (102초)"],
    ["15:43:26", "10", "225MB", "1분 43초 (103초)"],
], Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4),
   col_widths=[Inches(3.6), Inches(1.8), Inches(1.8), Inches(5.1)], font_size=15)

add_page_number(s, 7)


# ══════════════════════════════════════════════════════════════
# Slide 8 — Case 2 원인 & Before/After
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "Case 2 — CPU Spike: 원인 분석 & 조치", accent=GREEN)
add_bullets(s, [
    ("앱 내부 CPU 점유 로직이 레벨(1~10)에 따라 코어를 독점 — Lv=10에서는 5초간 코어 완전 점유", 0, True),
    ("약 1분 42초 주기로 CPU_MAX_OCCUPY(50%) 초과 → Watchdog이 SIGTERM으로 프로세스 종료", 0),
], Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.85), base_size=18)

add_bullets(s, [
    ("증거: \"Lv=10 → 5초간 코어 완전 점유\"는 아래 코드 로직과 실측 로그로 확인됨", 0, True, ACCENT),
], Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.4), base_size=17)

add_code_box(s, [
    "[agent-leak-sim.py — cpu_step(level) 핵심 로직]",
    "burn_sec = level * 0.5      # Lv=10 → 10 × 0.5 = 5.0초 (코어 점유 시간)",
    "_burn(burn_sec)              # 5.0초 동안 sqrt 연산 반복 (busy-loop, CPU 점유)",
    "actual = (CPU시간 증가량 / wall시간) × 100   # 실측 점유율(%)",
    "",
    "[실측 로그 — Lv=8(burn_sec=4.0s) 동일 로직 적용 결과]",
    "17:37:50 [INFO] [CPU] Occupy core for 4.0s (Level 8)",
    "17:37:50 [WARNING] [Watchdog] CPU usage spike detected: 99.7% > 50%",
    "  → wall ≈ 4.0s, actual ≈ 99.7%  (거의 100% 코어 점유 확인)",
    "  → 동일 busy-loop이 Lv=10(5.0s)에도 적용 → 5.0초간 ≈ 100% 점유",
], Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.55), font_size=12)

add_table(s, [
    ["구분", "설정", "결과"],
    ["Before", "CPU_MAX_OCCUPY = 50 %", "Lv=10 도달 시 99.7% > 50% → 즉시 WATCHDOG 종료"],
    ["After", "CPU_MAX_OCCUPY = 80~90 %", "더 높은 사용률 허용 → 안정적 동작 유지"],
], Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.05),
   col_widths=[Inches(2.0), Inches(4.0), Inches(6.3)], font_size=14)

add_bullets(s, [
    ("✅ 근본 해결: CPU를 독점하는 연산을 여러 작업으로 분산하거나, 중간에 sleep/yield를 삽입하여 다른 프로세스에 양보", 0, True, GREEN),
], Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), base_size=16)

add_page_number(s, 8)


# ══════════════════════════════════════════════════════════════
# Slide 9 — Case 3 현상 & 증거
# ══════════════════════════════════════════════════════════════
case_evidence_slide(
    9, "Case 3 — Deadlock: 현상 & 증거",
    info_bullets=[
        ("현상: 프로세스가 종료되지 않고 PID 유지, CPU/MEM 변화 없이 로그가 완전히 정지", 0, True),
        ("조건: MULTI_THREAD_ENABLE=true / MEMORY_LIMIT=9999 / CPU_MAX_OCCUPY=100", 0),
    ],
    code_lines=[
        "[agent_app.log — 마지막 로그]",
        "17:45:28 [INFO] [Thread-A] Acquired Lock-1, WAITING for Lock-2...",
        "17:45:28 [INFO] [Thread-B] Acquired Lock-2, WAITING for Lock-1...",
        "(이후 로그 없음 — 영원히 멈춤)",
        "",
        "[ps -ef | grep agent-leak-sim]",
        "PID 28648   TIME 0:55.09   Python agent-leak-sim.py",
        "  -> PID 존재(생존) + TIME 값 변화 없음(작업 정지) = Deadlock",
    ],
    extra_bullets=[
        ("A는 B가 가진 Lock-2를, B는 A가 가진 Lock-1을 서로 기다리며 영원히 대기 = Deadlock", 0),
    ],
)


# ══════════════════════════════════════════════════════════════
# Slide 10 — Case 3 원인 & 조치 (4대 조건 포함)
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "Case 3 — Deadlock: 원인 분석 & 조치", accent=GREEN)
add_table(s, [
    ["Deadlock 4대 조건", "우리 실습에서의 발생 양상"],
    ["상호 배제 (Mutual Exclusion)", "threading.Lock() — 한 번에 한 스레드만 획득 가능"],
    ["점유 대기 (Hold and Wait)", "Thread-A가 Lock-1을 쥔 채 Lock-2를 대기"],
    ["비선점 (No Preemption)", "Thread-B가 보유한 Lock-2를 A가 강제로 가져올 수 없음"],
    ["순환 대기 (Circular Wait)", "A→Lock-2 대기, B→Lock-1 대기 (서로 순환 대기)"],
], Inches(0.5), Inches(1.25), Inches(12.3), Inches(2.6),
   col_widths=[Inches(4.3), Inches(8.0)], font_size=15)

add_table(s, [
    ["구분", "설정", "결과"],
    ["Before", "MULTI_THREAD_ENABLE = true", "교착상태 발생 → PID 생존, 로그/CPU 정지(무응답)"],
    ["After", "MULTI_THREAD_ENABLE = false", "싱글스레드 → 자원 경쟁 없음 → 정상 동작 확인"],
], Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.4),
   col_widths=[Inches(2.0), Inches(4.0), Inches(6.3)], font_size=15)

add_bullets(s, [
    ("✅ 근본 해결: 락 획득 순서를 통일(모든 스레드가 Lock-1→Lock-2 순서로만 획득)하거나,", 0, True, GREEN),
    ("lock.acquire(timeout=...)으로 타임아웃을 설정하여 일정 시간 후 대기를 포기하도록 개선", 1, False, GREEN),
], Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.3), base_size=18)
add_page_number(s, 10)


# ══════════════════════════════════════════════════════════════
# Slide 11 — Q&A ① 메모리 / CPU
# ══════════════════════════════════════════════════════════════
qa_slide(11, "설명 자료 Q&A ① — 메모리 / CPU", [
    ("메모리 누수가 시스템에 미치는 영향은?",
     "메모리는 OS 전체가 공유하는 자원. 한 프로세스가 계속 독점하면 다른 프로세스가 메모리를 "
     "할당받지 못해 시스템이 느려지거나 OOM 상태가 되고, 심하면 OS의 OOM Killer가 강제 종료시킴"),
    ("MEMORY_LIMIT을 높이는 것이 근본 해결책인가?",
     "아니다 — 임시방편. 메모리 누수(해제 안 함) 자체는 그대로이므로, 한계치를 올려도 결국 "
     "그 한계마저 초과하여 다시 종료됨. 코드에서 사용 후 메모리를 해제해야 근본 해결"),
    ("CPU 과점유가 시스템 전체에 미치는 영향은?",
     "CPU는 여러 프로세스가 시분할(time-sharing)로 나눠 쓰는 자원. 하나가 독점하면 다른 "
     "프로세스들의 응답이 지연되고, 웹 서버라면 모든 요청이 느려져 타임아웃 발생"),
    ("Watchdog이 SIGTERM을 보내는 이유는?",
     "SIGTERM은 “정상적으로 마무리하고 종료하라”는 신호. 파일 닫기, 로그 저장 등 마무리 "
     "작업 후 종료되어 데이터 손실이 없음 (반면 SIGKILL은 즉시 강제 종료, 최후의 수단)"),
])


# ══════════════════════════════════════════════════════════════
# Slide 12 — Q&A ② Deadlock / 리포팅
# ══════════════════════════════════════════════════════════════
qa_slide(12, "설명 자료 Q&A ② — Deadlock / 리포팅", [
    ("Deadlock 프로세스를 어떻게 식별하나요?",
     "PID는 있는데 CPU/MEM 변화가 없고 로그도 멈췄다면 교착상태를 의심. top -H로 모든 스레드 "
     "상태가 대기(S, D)인지 확인하고, 마지막 로그에서 WAITING / BLOCKED 패턴을 찾음"),
    ("MULTI_THREAD_ENABLE=false가 근본 해결책인가요?",
     "아니다 — 멀티스레드를 끄면 Deadlock은 없어지지만 동시 처리 성능 이점도 사라짐. 근본 "
     "해결은 락 획득 순서 통일 또는 lock.acquire(timeout=...) 적용"),
    ("GitHub Issue 형태로 리포트를 작성하는 이유는?",
     "(1) 재현 방법 공유 — 같은 환경변수로 같은 문제 재현 가능, (2) Before/After 검증 기록 — "
     "조치 효과를 누구나 확인 가능, (3) 유사 장애 시 참고 자산 — “내 머릿속 지식”을 “팀 자산”으로 전환"),
], accent=GREEN)


# ══════════════════════════════════════════════════════════════
# Slide 13 — 제출 체크리스트
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "제출 체크리스트 — 필수 포함 항목 점검")
add_table(s, [
    ["필수 항목", "OOM (Case 1)", "CPU (Case 2)", "Deadlock (Case 3)"],
    ["1. 현상 설명 (Description)", "✅", "✅", "✅"],
    ["2. 로그 발췌 (monitor.log / agent_app.log)", "✅", "✅", "✅"],
    ["3. 앱 실행 로그 발췌", "✅", "✅", "✅"],
    ["4. ps / top 출력 캡처", "✅", "✅", "✅"],
    ["5. 근본 원인 분석 (Root Cause)", "✅", "✅", "✅"],
    ["6. Before & After 비교 (Workaround)", "✅", "✅", "✅"],
], Inches(0.5), Inches(1.25), Inches(12.3), Inches(3.6),
   col_widths=[Inches(5.5), Inches(2.27), Inches(2.27), Inches(2.26)], font_size=16)

add_bullets(s, [
    ("생성된 결과물:", 0, True),
    ("agent-leak-sim.py — 3가지 장애를 sudo 없이 재현하는 Python 시뮬레이터", 1),
    ("report-case1-memory-leak.md / report-case2-cpu-spike.md / report-case3-deadlock.md", 1),
    ("위 3개 보고서를 GitHub Issue로 등록 (linksat1/AI-SW-Basic) → 최종 제출", 1),
], Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0), base_size=17)
add_page_number(s, 13)


# ══════════════════════════════════════════════════════════════
# Slide 14 — 전체 흐름 요약 & 마무리
# ══════════════════════════════════════════════════════════════
s = new_slide()
add_title_bar(s, "전체 흐름 요약 & 마무리")
add_bullets(s, [
    ("①  agent-leak-sim.py 실행 → 케이스 선택 (1: Memory Leak / 2: CPU Spike / 3: Deadlock)", 0, True),
    ("②  각 케이스 실행 → agent_app.log + ps/top 캡처로 증거 수집", 0, True),
    ("③  로그 기반 원인 분석 → 환경변수 조정 후 Before/After 비교로 효과 검증", 0, True),
    ("④  GitHub Issue 3건 작성 (Description / Evidence / Root Cause / Workaround)", 0, True),
    ("⑤  최종 제출 (GitHub Repository 링크 또는 Issue 3건)", 0, True),
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.6), base_size=22)

box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4))
box.fill.solid()
box.fill.fore_color.rgb = LIGHTBLUE
box.line.fill.background()
tf = box.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "“장애 = 감으로 추측하는 것” → “장애 = 데이터로 증명하고 기록하는 것” 으로의 전환 완료"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER
add_page_number(s, 14)


# ── 저장 ────────────────────────────────────────────────────────
OUT = "/Users/cspag5955/AI-SW-Basic-1/b1-2/B1-2_미션_발표자료.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)}")
